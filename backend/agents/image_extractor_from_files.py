import os
import io
import base64
import asyncio

from openai import AzureOpenAI
from azure.core.credentials import AzureKeyCredential
from openpyxl import load_workbook
from openpyxl.drawing.image import Image
from PIL import Image as PILImage
from docx import Document
from openai import AzureOpenAI  
from azure.identity import DefaultAzureCredential, get_bearer_token_provider  
from pptx import Presentation
from backend.tools.LLM import MODELS, API_VERSION

from dotenv import load_dotenv
import requests
from urllib.parse import urlparse
import logging
import subprocess
import fitz  # PyMuPDF

####################################################
## set up logging
####################################################
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Load environment variables from .env file
load_dotenv(override=True)

output_results = []

def extract_images_from_pdf(pdf_path):
    images = []
    doc = fitz.open(pdf_path)
    for page_index in range(len(doc)):
        page = doc.load_page(page_index)
        for image_info in page.get_images(full=True):
            xref = image_info[0]
            base_image = doc.extract_image(xref)
            images.append(base_image["image"])
    return images


def extract_images_from_pptx(pptx_path):
    presentation = Presentation(pptx_path)
    images = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # 13 corresponds to a picture
                image = shape.image
                image_bytes = image.blob
                images.append(image_bytes)

    return images

def print_wrapped_text(text, max_length):
    lines = []
    words = text.split()
    current_line = ""
    for word in words:
        if len(current_line) + len(word) + 1 <= max_length:
            current_line += (word + " ")
        else:
            print(current_line.strip())
            current_line = word + " "
    if current_line:
        print(current_line.strip())
        lines.append(current_line.strip())
    return lines

def extract_images_from_docx(docx_path):
    logger.info(f"Extracting images from word doc")
    doc = Document(docx_path)
    images = []

    try:
        convert_office_file_to_pdf(docx_path, "temp.pdf")
        logger.info(f"Converting word doc to pdf")
    except Exception as e:
        logger.error(f"Error converting word doc to pdf: {e}")

    for rel in doc.part.rels.values():
        
        if "image" in rel.target_ref:
            # page_number = rel.target_part.page_number
            # logger.info(f"Page number: {page_number}")
            image = rel.target_part.blob
            images.append(image)
    return images

def convert_office_file_to_pdf(input_file: str, output_file: str) -> None:
    command = [
        "libreoffice",
        "--headless",
        "--convert-to",
        "pdf",
        input_file,
        "--outdir",
        os.path.dirname(output_file),
    ]
    subprocess.run(command, check=True)


def extract_images_from_xlsx(xlsx_path):
    workbook = load_workbook(xlsx_path)
    images = []

    for sheet in workbook.worksheets:
        for drawing in sheet._images:
            img = drawing._data()
            images.append(img)

    if not images:
        print("No images found in the workbook.")
    else:
        print(f"Found {len(images)} images in the workbook.")
 
    return images

def load_image_as_base64(image_path):
    with open(image_path, "rb") as image_file:
        file = image_file.read()
    return file



# https://github.com/Coding-Forge/Fabric/blob/main/images/admin-portal-settings.png

def upload_and_extract_images(file_path):
    _, file_extension = os.path.splitext(file_path)
    print(f"File path: {file_path}")
    if file_extension.lower() == '.docx':
        images = extract_images_from_docx(file_path)
    elif file_extension.lower() == '.xlsx':
        images = extract_images_from_xlsx(file_path)
    elif file_extension.lower() in ['.png', '.jpg', '.jpeg']:
        images = [load_image_as_base64(file_path)]
    elif file_extension.lower() == '.pptx':
        images = extract_images_from_pptx(file_path)
    elif file_extension.lower() == '.pdf':
        images = extract_images_from_pdf(file_path)
    else:
        raise ValueError("Unsupported file type")
    
    return images

async def receive_binary_files(binary_files: list[dict[str, str]]):
    results = []
    for file_data in binary_files:
        encoded = base64.b64encode(file_data['file']).decode('utf-8')
        results.append(f"data:application/octet-stream;base64,{encoded}")
    return results

async def workwithfile(filename:str):
    # Azure OpenAI Key Authentication

    try:
        images = upload_and_extract_images( filename )
    except Exception as e:
        print(f"Error loading workbook: {e}")
        exit()

    # ext = os.path.splitext(filename)[1]

    # Convert images to base64 strings
    # save all the images to a list
    image_urls = []

    for img in images:
        
        if isinstance(img, str) and img.startswith("data:image"):
            img_base64 = img.split(",")[1]
        else:
            img_base64 = base64.b64encode(img).decode('utf-8')

        image_urls.append(f"data:image/png;base64,{img_base64}")
        
        await workwithbase64encodedimages(image_urls)


async def workwithbase64encodedimages(image_urls:list):
    for image_url in image_urls:

        client = AzureOpenAI(
            azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
            api_key=os.getenv("AZURE_OPENAI_API_KEY"),
            api_version=API_VERSION.GPT4o.value,
        )

        if not client:
            print("Client not initialized")
            exit()

        response = client.chat.completions.create(
            model=MODELS.GPT4o.value,
            messages= [
                {
                    "role": "system",
                    "content": "You are an AI assistant that helps people find information. If you detect information in a table layout, provide feedback and return as a table using markdown.",
                },
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "Describe this picture"
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": image_url
                                # "url": [ image_urls ] #f"data:image/png;base64,{img_base64}"
                            }
                        }
                    ]
                }    
            ],
            max_tokens=2000
        )

        if response:
            result = response.choices[0].message 
            # output = print_wrapped_text(result.content, 150)
            output_results.append(result.content)


async def get_file(file_path: str):
    # await workwithfile(file_path)
    if file_path.startswith("http"):
        response = requests.get(file_path)
        if response.status_code == 200:

            parsed_url = urlparse(file_path)
            _, ext = os.path.splitext(parsed_url.path)
            temp_filename = "temp_downloaded_file" + (ext if ext else "")
            with open(temp_filename, "wb") as tmp:
                tmp.write(response.content)
            await workwithfile(temp_filename)

            os.remove(temp_filename)

        else:
            print("Failed to retrieve content from URL.")
    else:
        await workwithfile(file_path)

    if output_results:
        return output_results

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python image_extractor_from_files.py <file_path>")
        sys.exit(1)

    file_path = sys.argv[1]
    asyncio.run(get_file(file_path))